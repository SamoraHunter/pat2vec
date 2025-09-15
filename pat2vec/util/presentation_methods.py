import os
from typing import Dict, List

from pptx import Presentation
from pptx.util import Inches


def group_images_by_suffix(folder_path: str) -> Dict[str, List[str]]:
    """Groups image files in a folder based on their filename suffix.

    For example, 'image_A_client1.png' and 'image_B_client1.png' would both
    be grouped under the key 'client1'.

    Args:
        folder_path: The path to the folder containing the images.

    Returns:
        A dictionary where keys are suffixes and values are lists of image
        filenames.
    """
    image_groups = {}

    for file_name in os.listdir(folder_path):
        if file_name.endswith(".png") or file_name.endswith(".jpg") or file_name.endswith(".jpeg"):
            suffix = file_name.split("_")[-1].split(".")[0]
            image_groups.setdefault(suffix, []).append(file_name)

    return image_groups

def create_powerpoint_slides_client_idcode_groups(image_groups: Dict[str, List[str]], output_path: str) -> None:
    """Creates a PowerPoint presentation from grouped images.

    Each image from the `image_groups` dictionary is placed on a new, blank
    slide in the presentation.

    Args:
        image_groups: A dictionary where keys are group identifiers and values
            are lists of image filenames.
        output_path: The path where the output PowerPoint presentation will be saved.
    """
    presentation = Presentation()

    for suffix, images in image_groups.items():
        for image in images:
            image_path = os.path.join(os.path.dirname(output_path), image)
            slide_layout = presentation.slide_layouts[5]  # 5 corresponds to a blank slide layout
            slide = presentation.slides.add_slide(slide_layout)

            left = top = Inches(1)
            pic = slide.shapes.add_picture(image_path, left, top, height=Inches(5))

    presentation.save(output_path)
    print(f"PowerPoint file '{output_path}' created successfully.")

def create_powerpoint_from_images_group(folder_path: str) -> None:
    """Creates a PowerPoint presentation by grouping images in a folder.

    Args:
        folder_path: The path to the folder containing the images.
    """
    image_groups = group_images_by_suffix(folder_path)
    output_path = os.path.join(folder_path, 'output_presentation.pptx')
    create_powerpoint_slides_client_idcode_groups(image_groups, output_path)

#if __name__ == "__main__":
    #folder_path = '../plot_outputs_folder_piechart'
    #create_powerpoint_from_images_group(folder_path)


def create_powerpoint_slides(images: List[str], folder_path: str, output_path: str) -> None:
    """Creates a PowerPoint presentation with one image per slide.

    Args:
        images: A list of image filenames.
        folder_path: The directory where the image files are located.
        output_path: The path to save the generated PowerPoint file.
    """
    presentation = Presentation()

    for image in images:
        image_path = os.path.join(folder_path, image)
        slide_layout = presentation.slide_layouts[5]  # 5 corresponds to a blank slide layout
        slide = presentation.slides.add_slide(slide_layout)

        left = top = Inches(1)
        pic = slide.shapes.add_picture(image_path, left, top, height=Inches(5))

    presentation.save(output_path)
    print(f"PowerPoint file '{output_path}' created successfully.")

def create_powerpoint_from_images(folder_path: str) -> None:
    """
    Create a PowerPoint presentation from images in a specified folder.

    Parameters:
    - folder_path (str): The path to the folder containing images.
    """
    images = [file_name for file_name in os.listdir(folder_path)
              if file_name.endswith((".png", ".jpg", ".jpeg"))]

    output_path = os.path.join(folder_path, 'output_presentation.pptx')
    create_powerpoint_slides(images, folder_path, output_path)

# if __name__ == "__main__":
#     folder_path = '../plot_outputs_folder_piechart'
#     create_powerpoint_from_images(folder_path)
