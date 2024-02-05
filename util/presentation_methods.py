import os

from pptx import Presentation
from pptx.util import Inches


def group_images_by_suffix(folder_path):
    """
    Group images in a folder based on the suffix in their filenames.

    Parameters:
    - folder_path (str): The path to the folder containing images.

    Returns:
    - dict: A dictionary where keys are suffixes and values are lists of image filenames.
    """
    image_groups = {}

    for file_name in os.listdir(folder_path):
        if file_name.endswith(".png") or file_name.endswith(".jpg") or file_name.endswith(".jpeg"):
            suffix = file_name.split("_")[-1].split(".")[0]
            image_groups.setdefault(suffix, []).append(file_name)

    return image_groups

def create_powerpoint_slides_client_idcode_groups(image_groups, output_path):
    """
    Create a PowerPoint presentation with slides containing images from different groups.

    Parameters:
    - image_groups (dict): Dictionary where keys are suffixes and values are lists of image filenames.
    - output_path (str): The path for the output PowerPoint presentation.
    """
    presentation = Presentation()

    for suffix, images in image_groups.items():
        for image in images:
            image_path = os.path.join(output_path, image)
            slide_layout = presentation.slide_layouts[5]  # 5 corresponds to a blank slide layout
            slide = presentation.slides.add_slide(slide_layout)

            left = top = Inches(1)
            pic = slide.shapes.add_picture(image_path, left, top, height=Inches(5))

    presentation.save(output_path)
    print(f"PowerPoint file '{output_path}' created successfully.")

def create_powerpoint_from_images_group(folder_path):
    """
    Create a PowerPoint presentation from images in a specified folder.

    Parameters:
    - folder_path (str): The path to the folder containing images.
    """
    image_groups = group_images_by_suffix(folder_path)
    output_path = os.path.join(folder_path, 'output_presentation.pptx')
    create_powerpoint_slides_client_idcode_groups(image_groups, output_path)

#if __name__ == "__main__":
    #folder_path = '../plot_outputs_folder_piechart'
    #create_powerpoint_from_images_group(folder_path)


def create_powerpoint_slides(images, output_path):
    """
    Create a PowerPoint presentation with slides containing individual images.

    Parameters:
    - images (list): List of image filenames.
    - output_path (str): The path for the output PowerPoint presentation.
    """
    presentation = Presentation()

    for image in images:
        image_path = os.path.join(output_path, image)
        slide_layout = presentation.slide_layouts[5]  # 5 corresponds to a blank slide layout
        slide = presentation.slides.add_slide(slide_layout)

        left = top = Inches(1)
        pic = slide.shapes.add_picture(image_path, left, top, height=Inches(5))

    presentation.save(output_path)
    print(f"PowerPoint file '{output_path}' created successfully.")

def create_powerpoint_from_images(folder_path):
    """
    Create a PowerPoint presentation from images in a specified folder.

    Parameters:
    - folder_path (str): The path to the folder containing images.
    """
    images = [file_name for file_name in os.listdir(folder_path)
              if file_name.endswith((".png", ".jpg", ".jpeg"))]

    output_path = os.path.join(folder_path, 'output_presentation.pptx')
    create_powerpoint_slides(images, output_path)

# if __name__ == "__main__":
#     folder_path = '../plot_outputs_folder_piechart'
#     create_powerpoint_from_images(folder_path)
